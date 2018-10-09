import json
import pptx

from .utils import BaseParser


class Parser(BaseParser):
    """Extract text from pptx file using python-pptx
    """

    def extract(self, filename, **kwargs):
        """
        If `kwargs` has a keyword called `page` with value `True`,
        then the extraction will be done on a per page basis.
        The resultant text would be arranged in pages as it appears in
        original document.
        """
        presentation = pptx.Presentation(filename)
        page_extraction = kwargs.get('page', False)
        if page_extraction:
            text_runs = {}
            for index, slide in enumerate(presentation.slides):
                text_runs[index] = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs[index].append(run.text)
                text_runs[index] = '\n\n'.join(text_runs[index])
            return json.dumps(text_runs)
        else:
            text_runs = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
            return '\n\n'.join(text_runs)
